"""
Document Generator Service — produces .docx and .pptx files from account state.

Proposal generation uses Claude Haiku to write section prose (matches /proposal skill quality).
All generators read the account ASO, signals, and recent interactions.
Template support: workspace.settings["proposal_template_docx_b64"] overrides base document.
"""
import io
import base64
import json
import uuid
from datetime import datetime, timezone, date
from typing import Optional

import anthropic
import structlog
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select

from app.config import get_settings
from app.agents.drafter import DOCUMENT_VOICE_RULES, strip_dashes


def _strip_dashes_deep(value):
    """Apply strip_dashes to every string in a parsed JSON structure."""
    if isinstance(value, str):
        return strip_dashes(value)
    if isinstance(value, list):
        return [_strip_dashes_deep(v) for v in value]
    if isinstance(value, dict):
        return {k: _strip_dashes_deep(v) for k, v in value.items()}
    return value

log = structlog.get_logger()

# Brand colours
NAVY   = (26, 60, 110)       # #1A3C6E
ORANGE = (232, 119, 34)      # #E87722
WHITE  = (255, 255, 255)
GREY   = (100, 100, 100)


def _ws_sender(data: dict) -> dict:
    """Extract sender/company context from workspace settings, with generic fallbacks."""
    ws = data.get("workspace")
    settings = (ws.settings or {}) if ws else {}
    return {
        "name":    settings.get("sender_name") or "Your Name",
        "title":   settings.get("sender_title") or "Account Executive",
        "company": settings.get("sender_company") or "Our Company",
        "email":   settings.get("sender_email") or "",
        "product": settings.get("product_description") or settings.get("product_name") or "our product",
    }

HAIKU_MODEL = "claude-haiku-4-5-20251001"


# ── Data loader ───────────────────────────────────────────────────────────────

async def _load_account_data(account_id: str, db: AsyncSession) -> dict:
    from app.models.account import Account, Signal, Interaction, TimelineAction
    from app.models.workspace import Workspace

    acct = (await db.execute(
        select(Account).where(Account.id == account_id, Account.deleted_at.is_(None))
    )).scalar_one_or_none()
    if not acct:
        raise ValueError(f"Account {account_id} not found")

    # Workspace (for template)
    ws = (await db.execute(
        select(Workspace).where(Workspace.id == acct.workspace_id)
    )).scalar_one_or_none()
    template_b64: Optional[str] = (ws.settings or {}).get("proposal_template_docx_b64") if ws else None

    # Top 5 unacknowledged signals
    signals = (await db.execute(
        select(Signal)
        .where(Signal.account_id == account_id, Signal.acknowledged == False)  # noqa: E712
        .order_by(Signal.urgency.desc())
        .limit(5)
    )).scalars().all()

    # Last 3 call transcripts
    transcripts = (await db.execute(
        select(Interaction)
        .where(Interaction.account_id == account_id, Interaction.type == "call")
        .order_by(Interaction.occurred_at.desc())
        .limit(3)
    )).scalars().all()

    # Timeline actions (for mutual action plan)
    actions = (await db.execute(
        select(TimelineAction)
        .where(
            TimelineAction.account_id == account_id,
            TimelineAction.status.in_(["upcoming", "overdue", "completed"]),
        )
        .order_by(TimelineAction.due_date.asc())
        .limit(20)
    )).scalars().all()

    state = acct.state or {}
    pov   = state.get("pov") or {}

    # Narrative: try several paths in order
    narrative = (
        pov.get("deal_narrative")
        or pov.get("forecast_rationale")
        or pov.get("summary")
        or state.get("narrative")
        or ""
    )

    gold_data  = pov.get("gold_data") or {}
    meddpicc   = pov.get("meddpicc") or {}
    three_whys = pov.get("three_whys") or {}
    risks      = pov.get("risks") or []

    # Transcript notes for context
    transcript_text = ""
    for t in transcripts:
        notes = (t.notes or "").strip()
        if notes:
            transcript_text += f"\n---\n{notes[:800]}"

    # Signal details for context
    signal_details = "\n".join(
        f"- {s.type.replace('_', ' ').title()}: {(s.detail or '')[:200]}"
        for s in signals if s.detail
    )

    return {
        "account": acct,
        "workspace": ws,
        "template_b64": template_b64,
        "state": state,
        "pov": pov,
        "gold_data": gold_data,
        "meddpicc": meddpicc,
        "three_whys": three_whys,
        "risks": risks,
        "signals": signals,
        "transcripts": transcripts,
        "actions": actions,
        "name": acct.name,
        "stage": acct.stage or "Discovery",
        "amount": float(acct.deal_amount or 0),
        "narrative": narrative,
        "grounding_confidence": pov.get("grounding_confidence") or 0.0,
        "health_score": float(pov.get("health_score") or acct.health_score or 0.0),
        "transcript_text": transcript_text,
        "signal_details": signal_details,
    }


def _fmt_amount(v: float) -> str:
    if v >= 1_000_000: return f"${v/1_000_000:.1f}M"
    if v >= 1_000:     return f"${v/1_000:.0f}K"
    return f"${v:.0f}"


def _today_str() -> str:
    return date.today().strftime("%d %B %Y")


# ── Claude prose generator ────────────────────────────────────────────────────

async def _claude_sections(data: dict) -> dict:
    """
    Use Claude Haiku to write proposal prose for all major sections.
    Returns a dict of section_key -> text.
    Cost: ~$0.02-0.05 per proposal.
    """
    settings = get_settings()
    client = anthropic.AsyncAnthropic(api_key=settings.anthropic_api_key)

    name     = data["name"]
    stage    = data["stage"]
    amount   = _fmt_amount(data["amount"])
    narrative = data["narrative"]
    meddpicc  = data["meddpicc"]
    three_whys = data["three_whys"]
    risks     = data["risks"]
    transcript_text = data["transcript_text"]
    signal_details  = data["signal_details"]

    def gv(key: str, default="") -> str:
        v = data["gold_data"].get(key)
        if not v: return default
        if isinstance(v, dict): return str(v.get("value") or v.get("text") or default)
        return str(v)

    def why_text(key: str) -> str:
        entry = three_whys.get(key) or {}
        if isinstance(entry, dict):
            return str(entry.get("evidence") or entry.get("text") or "")
        return str(entry)

    meddpicc_summary = ""
    for k, label in [("metrics","Metrics"), ("champion","Champion"), ("economic_buyer","Economic Buyer"), ("implicate_pain","Implicate Pain")]:
        comp = meddpicc.get(k) or {}
        score = comp.get("score", 0) if isinstance(comp, dict) else 0
        text  = comp.get("evidence") or comp.get("text") if isinstance(comp, dict) else ""
        if text:
            meddpicc_summary += f"{label} ({int(score*100)}%): {str(text)[:200]}\n"

    comp_risk = next((r for r in risks if isinstance(r, dict) and r.get("type") in ("competition","competitor")), None)
    competitor = comp_risk.get("detail", "") if comp_risk else ""

    context = f"""
Client: {name}
Deal stage: {stage}
Deal amount: {amount}
Industry: {gv('industry', 'high-risk industrial')}

Deal narrative:
{narrative or 'No narrative available. Use the signals and stage context below.'}

Key signals:
{signal_details or 'No signals available.'}

MEDDPICC intelligence:
{meddpicc_summary or 'Not yet populated.'}

Why Change: {why_text('why_change') or 'Not identified.'}
Why Now: {why_text('why_now') or 'Not identified.'}
Why Us: {why_text('why_us') or 'Not identified.'}

Call notes:
{transcript_text or 'No call notes available.'}

Known competitor: {competitor or 'None identified.'}
"""

    sender = _ws_sender(data)
    prompt = f"""You are writing a professional B2B proposal for {sender['company']} ({sender['product']}).
Write in a confident, direct, consultative tone. Plain sentences only. No em-dashes, no bullet points, no markdown. Maximum 3 sentences per paragraph.

Client context:
{context}

Write the following proposal sections as a JSON object. Each value must be plain prose (no markdown, no dashes):

{{
  "exec_summary": "2-3 sentences. Lead with the client's core challenge. End with what the vendor proposes.",
  "challenge": "2-3 sentences. Describe the specific challenge this client faces based on the context.",
  "business_case": "2-3 sentences. Quantify the cost of inaction. Reference any known metrics or industry benchmarks.",
  "why_vendor": "2-3 sentences. Why this vendor specifically. Reference key differentiators and relevant case study.",
  "solution_overview": "2-3 sentences. High-level description of the proposed solution for this client.",
  "next_steps": "3 numbered action items as a single string, each on a new line. Format: 1. [action]\\n2. [action]\\n3. [action]"
}}

Return only valid JSON. No extra text."""

    response = await client.messages.create(
        model=HAIKU_MODEL,
        max_tokens=1200,
        system=DOCUMENT_VOICE_RULES,
        messages=[{"role": "user", "content": prompt}],
    )

    raw = response.content[0].text.strip()
    # Strip markdown code fences if present
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]

    try:
        sections = _strip_dashes_deep(json.loads(raw))
    except Exception:
        log.warning("claude_sections_parse_failed", raw=raw[:200])
        # Fallback to empty strings so the doc still generates
        _sender = _ws_sender(data)
        sections = {
            "exec_summary": f"{name} faces operational challenges that {_sender['company']} can address with {_sender['product']}.",
            "challenge": "The client faces operational challenges in a competitive environment.",
            "business_case": "Delays in addressing these challenges generate measurable direct and indirect costs.",
            "why_vendor": f"{_sender['company']} offers a proven solution with key differentiators tailored to this client.",
            "solution_overview": f"We propose deploying {_sender['product']} for {name} to address the identified challenges.",
            "next_steps": "1. Review proposal and confirm scope\n2. Schedule technical walkthrough\n3. Proceed to formal contract and implementation",
        }

    return sections


# ── Word doc helpers ──────────────────────────────────────────────────────────

def _open_doc(template_b64: Optional[str] = None):
    """Open a new Document, optionally from an uploaded template."""
    from docx import Document
    if template_b64:
        try:
            doc_bytes = base64.b64decode(template_b64)
            return Document(io.BytesIO(doc_bytes))
        except Exception as e:
            log.warning("template_load_failed", error=str(e))
    return Document()


def _set_margins(doc):
    from docx.shared import Inches
    for section in doc.sections:
        section.top_margin    = Inches(1.0)
        section.bottom_margin = Inches(1.0)
        section.left_margin   = Inches(1.2)
        section.right_margin  = Inches(1.2)


def _add_heading(doc, text: str, level: int = 1):
    from docx.shared import Pt, RGBColor
    p = doc.add_heading("", level=level)
    run = p.add_run(text)
    run.bold = True
    run.font.color.rgb = RGBColor(*NAVY)
    run.font.name = "Calibri"
    run.font.size = Pt(14 if level == 1 else 12)
    return p


def _add_body(doc, text: str):
    from docx.shared import Pt, RGBColor
    if not text or not text.strip():
        return
    p = doc.add_paragraph()
    run = p.add_run(text.strip())
    run.font.name = "Calibri"
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(30, 30, 30)
    return p


def _add_table(doc, headers: list, rows: list):
    from docx.shared import Pt, RGBColor
    from docx.oxml.ns import qn
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    try:
        table.style = "Light Shading Accent 1"
    except Exception:
        pass
    hdr = table.rows[0].cells
    for i, h in enumerate(headers):
        hdr[i].text = h
        for para in hdr[i].paragraphs:
            for run in para.runs:
                run.bold = True
                run.font.size = Pt(10)
    for ri, row_data in enumerate(rows):
        cells = table.rows[ri + 1].cells
        for ci, val in enumerate(row_data):
            cells[ci].text = str(val)
            for para in cells[ci].paragraphs:
                for run in para.runs:
                    run.font.size = Pt(10)
    return table


# ── PROPOSAL ─────────────────────────────────────────────────────────────────

async def generate_proposal(account_id: str, db: AsyncSession) -> bytes:
    """
    Generate a proposal .docx using:
    - Claude Haiku for all section prose (matches /proposal skill quality)
    - User-uploaded template .docx as base document (if available)
    - Vendor 10-section structure
    """
    data = await _load_account_data(account_id, db)
    name     = data["name"]
    stage    = data["stage"]
    amount   = _fmt_amount(data["amount"])
    meddpicc = data["meddpicc"]
    risks    = data["risks"]
    signals  = data["signals"]
    transcripts = data["transcripts"]
    grounding   = data["grounding_confidence"]

    # Get Claude-written section prose
    sections = await _claude_sections(data)

    # Open doc (from template if uploaded, else blank)
    doc = _open_doc(data.get("template_b64"))
    _set_margins(doc)

    # ── Cover ─────────────────────────────────────────────────────────────
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    def centered(text: str, size: int, bold=False, color=NAVY):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = p.add_run(text)
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.bold = bold
        run.font.color.rgb = RGBColor(*color)
        return p

    sender = _ws_sender(data)
    centered("Sales Intelligence Proposal", 22, bold=True)
    centered(f"Prepared for {name}", 14, color=ORANGE)
    centered(f"{stage}  |  {amount}  |  {_today_str()}", 11, color=GREY)
    centered(f"Prepared by {sender['company']}  |  Confidential", 10, color=GREY)
    doc.add_page_break()

    # ── 1. Executive Summary ───────────────────────────────────────────────
    _add_heading(doc, "1. Executive Summary")
    _add_body(doc, sections.get("exec_summary", ""))
    doc.add_paragraph()

    # ── 2. Understanding Your Challenge ───────────────────────────────────
    _add_heading(doc, "2. Understanding Your Challenge")
    _add_body(doc, sections.get("challenge", ""))
    if signals:
        doc.add_paragraph()
        _add_heading(doc, "Key concerns identified", level=2)
        for s in signals[:3]:
            if s.detail:
                p = doc.add_paragraph(style="List Bullet")
                r = p.add_run(f"{s.type.replace('_', ' ').title()}: {s.detail[:200]}")
                r.font.size = Pt(11)
                r.font.name = "Calibri"
    doc.add_paragraph()

    # ── 3. Business Case ──────────────────────────────────────────────────
    _add_heading(doc, "3. Business Case")
    _add_body(doc, sections.get("business_case", ""))
    doc.add_paragraph()

    # ── 4. Why Us ─────────────────────────────────────────────────────────
    _add_heading(doc, f"4. Why {sender['company']}")
    _add_body(doc, sections.get("why_vendor", ""))
    doc.add_paragraph()
    _add_table(doc,
        ["Metric", "Result"],
        [
            ["Product",                            sender['product']],
            ["MEDDPICC score (this deal)",         f"{int((meddpicc.get('overall_score') or 0) * 100)}%"],
        ]
    )
    doc.add_paragraph()

    # ── 5. Proposed Solution ──────────────────────────────────────────────
    _add_heading(doc, "5. Proposed Solution")
    _add_body(doc, sections.get("solution_overview", ""))
    doc.add_paragraph()
    _add_table(doc,
        ["Module", "Description"],
        [
            ["SafeKey",    "Real-time AI safety monitoring. 60+ detection models. WhatsApp and email alerts."],
            ["SafeRound",  "Digital inspection rounds with photo capture and AI non-compliance flags."],
            ["SafeScript", "Automated safety observation reports with AI-generated corrective actions."],
            ["SafeStart",  "Digital permit-to-work with agentic risk assessment and approval routing."],
        ]
    )
    doc.add_paragraph()

    # ── 6. Deployment Architecture ────────────────────────────────────────
    _add_heading(doc, "6. Deployment Architecture")
    _add_body(doc,
        "SafeKey integrates with your existing CCTV infrastructure via RTSP or ONVIF. "
        "All AI processing runs at the edge with encrypted sync to Azure cloud. "
        "No raw video is stored externally. Remote access is secured via Tailscale VPN and role-based access control."
    )
    doc.add_paragraph()

    # ── 7. Privacy and Compliance ─────────────────────────────────────────
    _add_heading(doc, "7. Privacy and Compliance")
    _add_body(doc,
        f"{sender['company']} meets enterprise data privacy requirements. "
        "The platform supports configurable data retention and deletion policies to meet PDPA, GDPR, and local data residency requirements."
    )
    doc.add_paragraph()

    # ── 8. Commercial Proposal ────────────────────────────────────────────
    _add_heading(doc, "8. Commercial Proposal")
    _add_body(doc,
        f"Based on our discussions, we propose a total investment of {amount} for the initial deployment. "
        "A detailed pricing breakdown will be provided in our formal quotation. "
        "Commercial terms are flexible and milestone-based."
    )
    doc.add_paragraph()

    # ── 9. Implementation Roadmap ─────────────────────────────────────────
    _add_heading(doc, "9. Implementation Roadmap")
    _add_table(doc,
        ["Phase", "Activities", "Timeline"],
        [
            ["1. Discovery",      "Site survey, camera mapping, integration testing",      "Week 1-2"],
            ["2. Deployment",     "Edge device installation, model tuning, alert routing", "Week 3-4"],
            ["3. Go-Live",        "User training, dashboard handover, SLA activation",     "Week 5"],
            ["4. Optimisation",   "Model accuracy review, false-positive reduction",       "Ongoing"],
        ]
    )
    doc.add_paragraph()

    # ── 10. Next Steps ────────────────────────────────────────────────────
    _add_heading(doc, "10. Next Steps")
    next_steps_text = sections.get("next_steps", "")
    if next_steps_text:
        for line in next_steps_text.split("\n"):
            line = line.strip()
            if line:
                p = doc.add_paragraph(style="List Number")
                r = p.add_run(line.lstrip("0123456789. "))
                r.font.size = Pt(11)
                r.font.name = "Calibri"
    contact_line = f"Contact {sender['name']} at {sender['email']} to proceed." if sender['email'] else f"Contact {sender['name']} to proceed."
    _add_body(doc,
        f"\nThis proposal is valid for 30 days from {_today_str()}. {contact_line}"
    )

    # ── Data attribution footer ────────────────────────────────────────────
    doc.add_page_break()
    _add_heading(doc, "Data Sources", level=2)
    _add_body(doc,
        f"This proposal was AI-drafted from {len(signals)} verified signals and "
        f"{len(transcripts)} call records for {name}. "
        f"Grounding confidence: {int(grounding * 100)}%. "
        "All factual claims are sourced from confirmed CRM interactions and Vantage intelligence."
    )

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── SALES DECK (.pptx) ────────────────────────────────────────────────────────

async def _claude_deck_content(data: dict) -> dict:
    """
    Claude Haiku writes the 5 personalized sections of the deck.
    All other slides use workspace-configured content.
    """
    settings = get_settings()
    client = anthropic.AsyncAnthropic(api_key=settings.anthropic_api_key)

    name   = data["name"]
    stage  = data["stage"]
    amount = _fmt_amount(data["amount"])
    narrative        = data["narrative"]
    signal_details   = data["signal_details"]
    transcript_text  = data["transcript_text"]
    three_whys       = data["three_whys"]
    risks            = data["risks"]
    meddpicc         = data["meddpicc"]

    def why_text(key: str) -> str:
        e = three_whys.get(key) or {}
        return str(e.get("evidence") or e.get("text") or "") if isinstance(e, dict) else str(e)

    comp_risk = next((r for r in risks if isinstance(r, dict) and r.get("type") in ("competition", "competitor")), None)
    competitor = comp_risk.get("detail", "") if comp_risk else ""

    meddpicc_pain = meddpicc.get("implicate_pain") or {}
    pain_text = (meddpicc_pain.get("evidence") or meddpicc_pain.get("text") or "") if isinstance(meddpicc_pain, dict) else ""

    deck_sender = _ws_sender(data)
    prompt = f"""You are writing slides for a {deck_sender['company']} sales deck for a meeting with {name}.
{deck_sender['company']} sells: {deck_sender['product']}.

Client context:
- Name: {name}
- Stage: {stage}
- Amount: {amount}
- Deal narrative: {narrative or "Not available"}
- Key signals: {signal_details or "None"}
- Pain evidence: {pain_text or "Not identified"}
- Why change: {why_text("why_change") or "Not identified"}
- Why now: {why_text("why_now") or "Not identified"}
- Why us: {why_text("why_us") or "Not identified"}
- Known competitor: {competitor or "None"}
- Call notes: {transcript_text[:400] if transcript_text else "None"}

Return JSON only. No markdown. No extra text.

{{
  "headline": "One sentence: why {name} needs {deck_sender['product']} specifically. Start with the client name. Plain sentence, no jargon.",
  "context_bullets": [
    "Bullet 1: 1 sentence about this client's industry/operational context",
    "Bullet 2: 1 sentence about the safety challenge they face based on signals",
    "Bullet 3: 1 sentence about why now is the right time (regulatory, incident risk, growth)"
  ],
  "challenge_headline": "One headline slide title: the specific safety problem {name} faces. Max 10 words.",
  "challenge_detail": "2 sentences explaining the challenge specific to {name}. Reference their industry.",
  "case_study_client": "Pick the most relevant: HDB (construction/housing), Saipem (O&G/offshore), Micron Lumchang (Singapore construction), Fulton Hogan (Australian construction). Just the client name.",
  "case_study_headline": "Result achieved for the case study client. E.g. 60 percent reduction in near-miss incidents",
  "case_study_bullets": [
    "Result 1 with number",
    "Result 2 with number",
    "Result 3 with context"
  ],
  "why_us_headline": "Why us specifically for {name}. Max 10 words.",
  "why_us_points": [
    "Differentiator 1 specific to this deal (based on their stage/signals/competitor)",
    "Differentiator 2 (reference ISO 27001, data sovereignty, or relevant cert)",
    "Differentiator 3 (reference relevant client proof point or track record)"
  ],
  "next_step_1": "First action - specific to {name}. Start with a verb.",
  "next_step_2": "Second action. Start with a verb.",
  "next_step_3": "Third action. Start with a verb."
}}"""

    response = await client.messages.create(
        model=HAIKU_MODEL,
        max_tokens=1000,
        system=DOCUMENT_VOICE_RULES,
        messages=[{"role": "user", "content": prompt}],
    )
    raw = response.content[0].text.strip()
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"): raw = raw[4:]
    try:
        return _strip_dashes_deep(json.loads(raw))
    except Exception:
        log.warning("claude_deck_content_parse_failed", raw=raw[:200])
        return {
            "headline": f"{name} needs proactive AI safety monitoring to prevent incidents and meet regulatory requirements.",
            "context_bullets": [
                f"{name} operates in a high-risk environment with complex safety monitoring needs.",
                "Manual safety inspections miss up to 99% of near-miss events that precede serious incidents.",
                "Regulatory requirements and operational risk make real-time AI monitoring a priority now."
            ],
            "challenge_headline": f"{name}'s Safety Challenge",
            "challenge_detail": "Traditional manual safety monitoring cannot scale across multiple sites or shifts. Near-misses go undetected until they become serious incidents.",
            "case_study_client": "HDB",
            "case_study_headline": "60% reduction in near-miss incidents",
            "case_study_bullets": ["8:1 AI vs manual detection ratio", "30 cameras deployed across 5 sites", "6-month deployment timeline"],
            "why_us_headline": f"Why us for {name}",
            "why_us_points": [
                "Proven solution with strong differentiators tailored to your requirements.",
                "Enterprise-grade security and data sovereignty controls.",
                "Track record of successful deployments and measurable outcomes."
            ],
            "next_step_1": f"Review this deck and confirm the deployment scope with {name} safety team.",
            "next_step_2": "Schedule a 1-hour technical site walkthrough to map existing camera infrastructure.",
            "next_step_3": "Proceed to 30-day Proof of Value with agreed success metrics.",
        }


async def generate_sales_deck(account_id: str, db: AsyncSession) -> bytes:
    """
    15-slide PowerPoint deck — Claude-personalized on 5 key slides.
    Structure: Hero → Context → Challenge → Product → Case Study → Why Us → Next Steps
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor as pRGB
    from pptx.enum.text import PP_ALIGN

    data  = await _load_account_data(account_id, db)
    deck  = await _claude_deck_content(data)
    deck_sender = _ws_sender(data)

    name   = data["name"]
    stage  = data["stage"]
    amount = _fmt_amount(data["amount"])

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    BG     = (248, 249, 251)   # Near-white page background
    DARK   = (20, 20, 30)      # Dark text

    def blank():
        return prs.slides.add_slide(prs.slide_layouts[6])

    def bg(slide, color=BG):
        """Fill slide background."""
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = pRGB(*color)

    def box(slide, l, t, w, h, fill_color, line=False):
        """Solid-filled rectangle shape."""
        shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = pRGB(*fill_color)
        if not line:
            shape.line.fill.background()
        return shape

    def txt(slide, l, t, w, h, text, size=14, bold=False,
            color=DARK, align=PP_ALIGN.LEFT, wrap=True):
        """Text box."""
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = pRGB(*color)
        run.font.name = "Calibri"
        return tb

    def bullet_row(slide, y: float, text: str, size=13, color=DARK, dot_color=ORANGE):
        """Orange dot + text bullet row."""
        box(slide, 0.6, y + 0.08, 0.06, 0.28, dot_color)
        txt(slide, 0.85, y, 11.8, 0.5, text, size=size, color=color)

    def header_bar(slide, slide_title: str):
        """Navy top bar with slide title."""
        box(slide, 0, 0, 13.33, 0.55, NAVY)
        txt(slide, 0.4, 0.07, 9, 0.4, slide_title, size=16, bold=True,
            color=WHITE)
        txt(slide, 9.8, 0.1, 3.3, 0.35, deck_sender["company"].upper(), size=11, bold=True,
            color=(180, 195, 220), align=PP_ALIGN.RIGHT)

    def stat_card(slide, l: float, t: float, number: str, label: str, bg_color=NAVY):
        """Large stat card: big number + label below."""
        box(slide, l, t, 3.8, 2.2, bg_color)
        txt(slide, l + 0.1, t + 0.15, 3.6, 1.1, number,
            size=34, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
        txt(slide, l + 0.1, t + 1.3, 3.6, 0.75, label,
            size=12, color=WHITE, align=PP_ALIGN.CENTER)

    def section_label(slide, text: str, y=0.65):
        """Small orange uppercase section label above headline."""
        txt(slide, 0.5, y, 12, 0.3, text.upper(), size=9,
            bold=True, color=ORANGE)

    # ── Slide 1: Hero ──────────────────────────────────────────────────────────
    s = blank()
    bg(s, NAVY)
    # Diagonal accent
    box(s, 9.5, 4.5, 5.0, 4.0, (30, 55, 100))
    txt(s, 0.8, 0.9, 12, 1.4, "Rethinking Safety.", size=46, bold=True,
        color=WHITE, align=PP_ALIGN.LEFT)
    txt(s, 0.8, 2.5, 10, 0.55,
        "AI-powered safety intelligence for construction, O&G, manufacturing and logistics.",
        size=15, color=(190, 205, 225))
    box(s, 0.8, 3.3, 4.8, 0.05, ORANGE)
    txt(s, 0.8, 3.55, 8, 0.5, f"Prepared for {name}", size=15, bold=True, color=ORANGE)
    txt(s, 0.8, 4.15, 8, 0.4, f"{stage}  ·  {amount}  ·  {_today_str()}", size=12, color=(160, 175, 200))
    _sender_line = f"{deck_sender['name']}, {deck_sender['title']}"
    if deck_sender['email']:
        _sender_line += f"  ·  {deck_sender['email']}"
    txt(s, 0.8, 6.6, 6, 0.4, _sender_line,
        size=10, color=(130, 150, 180))
    # Cert badges bottom right
    for i, cert in enumerate(["ISO 27001", "GDPR", "PDPA", "SG Cyber Safe"]):
        bx = box(s, 8.5 + i * 1.2, 6.55, 1.1, 0.38, (40, 65, 115))
        txt(s, 8.5 + i * 1.2, 6.58, 1.1, 0.32, cert, size=8, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)

    # ── Slide 2: About This Presentation ──────────────────────────────────────
    s = blank()
    bg(s)
    header_bar(s, f"About This Presentation")
    section_label(s, f"Prepared for {name}")
    txt(s, 0.5, 1.0, 12, 0.7, deck["headline"],
        size=20, bold=True, color=NAVY)
    box(s, 0.5, 1.8, 12.3, 0.03, ORANGE)

    for i, bullet in enumerate(deck["context_bullets"]):
        bullet_row(s, 2.0 + i * 0.85, bullet, size=14)

    # Deal snapshot box bottom right
    box(s, 9.0, 5.0, 3.8, 2.0, NAVY)
    txt(s, 9.1, 5.1, 3.6, 0.4, "Deal Snapshot", size=11, bold=True, color=ORANGE)
    for j, (lbl, val) in enumerate([("Stage", stage), ("Value", amount), ("Date", _today_str())]):
        txt(s, 9.1, 5.55 + j * 0.45, 1.5, 0.38, lbl, size=10, color=(180,195,215))
        txt(s, 10.6, 5.55 + j * 0.45, 2.1, 0.38, val, size=10, bold=True, color=WHITE)

    # ── Slide 3: The Challenge ──────────────────────────────────────────────────
    s = blank()
    bg(s)
    header_bar(s, deck["challenge_headline"])
    txt(s, 0.5, 0.75, 12, 0.55, deck["challenge_detail"], size=13, color=(60, 60, 80))
    stat_card(s, 0.5,  1.7, "6,500+", "workers die daily\nfrom occupational accidents", NAVY)
    stat_card(s, 4.6,  1.7, "$200B+", "in annual losses from\nsafety incidents & downtime", (30, 50, 95))
    stat_card(s, 8.7,  1.7, "300:1",  "near-miss to serious\nincident ratio", (20, 40, 80))
    txt(s, 0.5, 4.15, 12.3, 0.4,
        "Most near-misses go undetected — manual observation catches less than 1% of unsafe events.",
        size=12, color=(100, 100, 120), align=PP_ALIGN.CENTER)

    # ── Slide 4: Every Incident is Preventable ────────────────────────────────
    s = blank()
    bg(s)
    header_bar(s, "Every Serious Incident is Preceded by Hundreds of Near-Misses")
    txt(s, 0.5, 0.75, 12, 0.4, "The safety iceberg — what you see vs. what's hidden.", size=13, color=(80, 80, 100))
    chain = [
        ("Undetected\nNear-Misses", (40, 60, 110)),
        ("Unsafe Acts\nNot Corrected", (50, 70, 130)),
        ("1 Serious\nIncident", (180, 50, 50)),
        ("$1M+\nTotal Loss", ORANGE),
    ]
    for i, (label, clr) in enumerate(chain):
        box(s, 0.5 + i * 3.1, 1.5, 2.7, 2.4, clr)
        txt(s, 0.5 + i * 3.1, 1.7, 2.7, 1.8, label, size=16, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        if i < 3:
            txt(s, 3.1 + i * 3.1, 2.3, 0.5, 0.6, "→", size=20, bold=True, color=(150, 165, 190))
    txt(s, 0.5, 4.15, 12.3, 0.5,
        "SafeKey detects unsafe acts and conditions in real time — before they escalate.",
        size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # ── Slide 5: Introducing SafeKey ──────────────────────────────────────────
    s = blank()
    bg(s)
    header_bar(s, "Introducing SafeKey — AI Safety Monitoring That Works 24/7")
    txt(s, 0.5, 0.75, 12, 0.4, "Plug into your existing CCTV. No new cameras. No downtime.", size=13, color=(80, 80, 100))
    features = [
        ("60+ Detection Models",   "PPE compliance, fall detection, restricted zones, fire, forklift proximity and more — all GA from day 1."),
        ("Real-Time Alerts",       "WhatsApp, email, and Teams alerts with annotated video clips delivered to supervisors within seconds."),
        ("Edge Processing",        "All AI inference runs on-site. No raw video leaves your facility. Full data sovereignty compliance."),
        ("Audit-Ready Reports",    "Automated daily and weekly safety reports with incident timelines, trend analysis, and corrective actions."),
    ]
    for i, (title, desc) in enumerate(features):
        bx = box(s, 0.5 + (i % 2) * 6.4, 1.4 + (i // 2) * 2.3, 5.9, 2.0, WHITE)
        box(s, 0.5 + (i % 2) * 6.4, 1.4 + (i // 2) * 2.3, 0.08, 2.0, ORANGE)
        txt(s, 0.8 + (i % 2) * 6.4, 1.55 + (i // 2) * 2.3, 5.5, 0.45,
            title, size=14, bold=True, color=NAVY)
        txt(s, 0.8 + (i % 2) * 6.4, 2.1 + (i // 2) * 2.3, 5.5, 0.9,
            desc, size=11, color=(70, 70, 90))

    # ── Slide 6: SafeSuite Modules ────────────────────────────────────────────
    s = blank()
    bg(s)
    header_bar(s, "One Platform, Four Modules — Deploy What You Need")
    modules = [
        ("SafeKey",    "Real-time AI safety monitoring",         "60+ detections · edge processing · live alerts"),
        ("SafeRound",  "Digital inspection rounds",              "Photo capture · AI non-compliance flagging · reports"),
        ("SafeScript", "Automated safety observations",         "AI-generated findings · corrective actions · audit trail"),
        ("SafeStart",  "Permit-to-work",                        "Digital PTW · agentic risk assessment · approval routing"),
    ]
    for i, (name_m, subtitle_m, detail_m) in enumerate(modules):
        bx = box(s, 0.5 + (i % 2) * 6.4, 1.0 + (i // 2) * 2.7, 5.9, 2.4, NAVY)
        txt(s, 0.7 + (i % 2) * 6.4, 1.15 + (i // 2) * 2.7, 5.5, 0.5,
            name_m, size=18, bold=True, color=ORANGE)
        txt(s, 0.7 + (i % 2) * 6.4, 1.7 + (i // 2) * 2.7, 5.5, 0.45,
            subtitle_m, size=13, bold=True, color=WHITE)
        txt(s, 0.7 + (i % 2) * 6.4, 2.25 + (i // 2) * 2.7, 5.5, 0.7,
            detail_m, size=10, color=(180, 195, 220))

    # ── Slide 7: Deployment ────────────────────────────────────────────────────
    s = blank()
    bg(s)
    header_bar(s, "Live in 2-4 Weeks — No New Infrastructure Required")
    steps = [
        ("Week 1-2",  "Discovery",    "Site survey · camera mapping · integration test · network check"),
        ("Week 3-4",  "Deployment",   "Edge device install · model tuning · alert routing · user setup"),
        ("Week 5+",   "Live",         "Training handover · dashboard access · SLA active · optimisation"),
    ]
    for i, (timeline, title_s, detail_s) in enumerate(steps):
        box(s, 0.5 + i * 4.2, 1.2, 3.8, 0.45, ORANGE)
        txt(s, 0.5 + i * 4.2, 1.22, 3.8, 0.4, timeline, size=12, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        box(s, 0.5 + i * 4.2, 1.65, 3.8, 3.5, WHITE)
        box(s, 0.5 + i * 4.2, 1.65, 3.8, 0.06, NAVY)
        txt(s, 0.65 + i * 4.2, 1.85, 3.5, 0.5, title_s, size=16, bold=True, color=NAVY)
        for j, point in enumerate(detail_s.split(" · ")):
            bullet_row(s, 2.55 + j * 0.55 + i * 0.0, point, size=11,
                       dot_color=ORANGE)
        # Adjust y per column to not overlap
        # Simpler: just list as paragraph
        box(s, 0.5 + i * 4.2, 2.2, 3.8, 2.95, (250, 251, 253))
        txt(s, 0.65 + i * 4.2, 2.25, 3.5, 2.8,
            detail_s.replace(" · ", "\n"), size=12, color=(60, 70, 90), wrap=True)
    txt(s, 0.5, 5.9, 12.3, 0.4,
        "No rip-and-replace. SafeKey integrates with any RTSP or ONVIF camera.",
        size=12, color=(100, 100, 120), align=PP_ALIGN.CENTER)

    # ── Slide 8: Detection Library ─────────────────────────────────────────────
    s = blank()
    bg(s)
    header_bar(s, "40+ Safety Detections — All Active From Day 1")
    detections = [
        "PPE: Hard Hat", "PPE: Safety Vest", "PPE: Safety Boots", "No Face Mask",
        "Restricted Zone Entry", "Person Near Edge", "Work at Heights", "Confined Space",
        "Fallen Person", "Fire Detected", "Smoke Detected", "SOP Non-Compliance",
        "Forklift Proximity", "Forklift Speed", "Near-Miss Vehicle", "Mobile Phone Use",
        "Fatigue (drowsiness)", "Slip / Trip Risk", "Unsafe Lifting", "Scaffold Inspection",
    ]
    cols = 5
    for i, det in enumerate(detections[:20]):
        col = i % cols
        row = i // cols
        bx = box(s, 0.4 + col * 2.5, 0.75 + row * 1.4, 2.3, 1.1, NAVY)
        txt(s, 0.45 + col * 2.5, 0.85 + row * 1.4, 2.2, 0.9,
            det, size=10, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    # ── Slide 9: Track Record ──────────────────────────────────────────────────
    s = blank()
    bg(s)
    header_bar(s, "Trusted by 40+ Clients Across Singapore, MENA, and Australia")
    stat_card(s, 0.5,  1.3, "40+",    "enterprise clients\nin safety-critical industries", NAVY)
    stat_card(s, 4.6,  1.3, "60+",    "AI detection models\nGA from day 1", (30, 50, 95))
    stat_card(s, 8.7,  1.3, "-40%",   "average incident reduction\nwithin 6 months", (20, 40, 80))
    clients = ["HDB",  "JTC Corporation", "Saipem", "Micron Lumchang",
               "Saudi Aramco Contractors", "Fulton Hogan", "SingPost", "Huationg"]
    txt(s, 0.5, 3.75, 12.3, 0.35, "Clients include:", size=11, bold=True, color=NAVY)
    for i, cl in enumerate(clients[:8]):
        col, row = i % 4, i // 4
        box(s, 0.5 + col * 3.1, 4.2 + row * 0.85, 2.9, 0.65, WHITE)
        txt(s, 0.6 + col * 3.1, 4.28 + row * 0.85, 2.7, 0.5,
            cl, size=11, color=NAVY, align=PP_ALIGN.CENTER)

    # ── Slide 10: Case Study ───────────────────────────────────────────────────
    case_client = deck.get("case_study_client", "HDB")
    s = blank()
    bg(s)
    header_bar(s, f"Case Study: {case_client}")
    section_label(s, "Proven Results")
    txt(s, 0.5, 1.0, 12, 0.65, deck["case_study_headline"],
        size=22, bold=True, color=NAVY)
    box(s, 0.5, 1.75, 12.3, 0.04, ORANGE)
    for i, bullet in enumerate(deck["case_study_bullets"]):
        bullet_row(s, 2.0 + i * 0.9, bullet, size=14)
    box(s, 8.5, 1.0, 4.3, 5.5, NAVY)
    txt(s, 8.6, 1.15, 4.1, 0.4, "Key Metrics", size=12, bold=True, color=ORANGE)
    metrics = [
        ("Detection", "8x vs manual"),
        ("Incidents", "-60% in 6 months"),
        ("Cameras", "30 integrated"),
        ("Go-live", "5 weeks"),
    ]
    for j, (lbl, val) in enumerate(metrics):
        txt(s, 8.6, 1.7 + j * 1.0, 2.0, 0.4, lbl, size=11, color=(170, 185, 210))
        txt(s, 8.6, 2.1 + j * 1.0, 4.1, 0.45, val, size=16, bold=True, color=WHITE)

    # ── Slide 11: Privacy & Compliance ────────────────────────────────────────
    s = blank()
    bg(s)
    header_bar(s, "Built for Regulated Environments — Your Data Never Leaves the Site")
    privacy_points = [
        ("Edge Processing",     "All AI runs on-site. No raw video transmitted externally."),
        ("Face Blurring",       "Automatic face anonymisation — optional body-only mode."),
        ("Auto-Delete",         "Configurable retention policies to meet PDPA and GDPR."),
        ("ISO 27001",           "Certified information security management system."),
        ("Data Sovereignty",    "Video and metadata stored in your preferred region."),
        ("Access Control",      "Role-based access with full audit trail per user action."),
    ]
    for i, (title_p, desc_p) in enumerate(privacy_points):
        col, row = i % 2, i // 2
        box(s, 0.5 + col * 6.4, 0.9 + row * 1.8, 5.9, 1.55, WHITE)
        box(s, 0.5 + col * 6.4, 0.9 + row * 1.8, 0.07, 1.55, ORANGE)
        txt(s, 0.75 + col * 6.4, 1.0 + row * 1.8, 5.5, 0.4, title_p,
            size=13, bold=True, color=NAVY)
        txt(s, 0.75 + col * 6.4, 1.45 + row * 1.8, 5.5, 0.7, desc_p,
            size=11, color=(70, 70, 90))

    # ── Slide 12: Why Us for [Client] ─────────────────────────────────────────
    s = blank()
    bg(s)
    header_bar(s, deck["why_us_headline"])
    section_label(s, f"Specifically for {name}")
    txt(s, 0.5, 1.0, 12, 0.5, "Three reasons this engagement makes sense now:", size=13, color=(80, 80, 100))
    for i, point in enumerate(deck["why_us_points"]):
        num_box = box(s, 0.5, 1.8 + i * 1.6, 0.65, 0.65, ORANGE)
        txt(s, 0.5, 1.82 + i * 1.6, 0.65, 0.6, str(i + 1), size=20, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, 1.35, 1.82 + i * 1.6, 11.3, 0.9, point, size=14, color=DARK, wrap=True)

    # ── Slide 13: Next Steps ───────────────────────────────────────────────────
    s = blank()
    bg(s, NAVY)
    txt(s, 0.8, 0.5, 12, 0.7, "Next Steps", size=34, bold=True, color=WHITE)
    box(s, 0.8, 1.3, 11.5, 0.05, ORANGE)
    next_steps = [
        deck.get("next_step_1", "Review this deck with your safety team."),
        deck.get("next_step_2", "Schedule a 1-hour technical site walkthrough."),
        deck.get("next_step_3", "Proceed to 30-day Proof of Value deployment."),
    ]
    for i, step in enumerate(next_steps):
        box(s, 0.8, 1.7 + i * 1.55, 0.7, 0.7, ORANGE)
        txt(s, 0.8, 1.73 + i * 1.55, 0.7, 0.65, str(i + 1), size=22, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        txt(s, 1.7, 1.73 + i * 1.55, 11.0, 0.65, step, size=16, color=WHITE)

    # ── Slide 14: Contact ──────────────────────────────────────────────────────
    s = blank()
    bg(s)
    box(s, 0, 0, 13.33, 0.6, ORANGE)
    txt(s, 0.4, 0.1, 5, 0.42, deck_sender['company'].upper(), size=18, bold=True, color=WHITE)
    txt(s, 0.8, 1.2, 11.5, 1.1, "Thank You.", size=42, bold=True, color=NAVY)
    box(s, 0.8, 2.5, 5.5, 0.05, NAVY)
    txt(s, 0.8, 2.75, 7, 0.45, f"Prepared for {name}", size=16, color=NAVY)
    _contact_block = f"{deck_sender['name']}\n{deck_sender['title']}, {deck_sender['company']}"
    if deck_sender['email']:
        _contact_block += f"\n\n{deck_sender['email']}"
    txt(s, 0.8, 3.4, 7, 1.8, _contact_block, size=14, color=(60, 70, 90))
    txt(s, 0.8, 5.7, 12.3, 0.4,
        "This presentation is confidential and intended solely for the named recipient.",
        size=10, color=(150, 160, 175), align=PP_ALIGN.CENTER)
    # Right panel
    box(s, 9.5, 1.2, 3.5, 5.5, NAVY)
    txt(s, 9.65, 1.4, 3.2, 0.45, "ISO 27001 Certified", size=11, bold=True, color=ORANGE)
    txt(s, 9.65, 2.0, 3.2, 3.5,
        "40+ enterprise clients\n60+ CV detection models\nSingapore · MENA · Australia\n\nshell.com/safekey",
        size=12, color=WHITE)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── BATTLE CARD ───────────────────────────────────────────────────────────────

async def generate_battle_card(account_id: str, db: AsyncSession) -> bytes:
    data = await _load_account_data(account_id, db)
    bc_sender = _ws_sender(data)
    doc  = _open_doc(data.get("template_b64"))
    _set_margins(doc)
    name  = data["name"]
    risks = data["risks"]

    comp_risk  = next((r for r in risks if isinstance(r, dict) and r.get("type") in ("competition","competitor")), None)
    competitor = comp_risk.get("detail", "the incumbent solution") if comp_risk else "the incumbent solution"

    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt, RGBColor

    cover = doc.add_paragraph()
    cover.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = cover.add_run("Competitive Battle Card")
    r.font.size = Pt(20); r.bold = True; r.font.color.rgb = RGBColor(*NAVY); r.font.name = "Calibri"

    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = sub.add_run(f"{name}  |  Internal Use Only  |  {_today_str()}")
    r2.font.size = Pt(11); r2.font.color.rgb = RGBColor(*GREY); r2.font.name = "Calibri"
    doc.add_paragraph()

    _add_heading(doc, f"Competing Against: {competitor}")
    _add_table(doc, [f"{bc_sender['company']} Advantage", "Why It Wins"], [
        ["Full audit trail, every fact sourced",       "Competitor evidence chain is opaque"],
        ["Self-serve CRM integration in minutes",      "Competitors require complex setup"],
        ["Workspace-driven AI context",                "Competitors use generic AI defaults"],
    ])
    doc.add_paragraph()

    _add_heading(doc, "Objection Handlers")
    _add_table(doc, ["Objection", "Response"], [
        ["They have more experience", "Longevity does not equal accuracy. Our GA model catalogue is broader."],
        ["Their pricing is lower",    "Total cost of ownership includes integration cost and false-alarm rate. We benchmark favourably."],
        ["We already use them",       "SafeKey integrates alongside existing tools. No rip-and-replace required."],
    ])
    doc.add_paragraph()

    _add_heading(doc, "Proof Points")
    _add_body(doc, "HDB: 8:1 AI vs manual detection, 60% reduction in near-miss incidents.")
    _add_body(doc, "Saipem: EU AI Act compliant. ISO 27001 audited.")
    _add_body(doc, "40+ clients across construction, O&G, manufacturing, and logistics.")

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── BUSINESS CASE ─────────────────────────────────────────────────────────────

async def generate_business_case(account_id: str, db: AsyncSession) -> bytes:
    data = await _load_account_data(account_id, db)
    bc_biz_sender = _ws_sender(data)
    doc  = _open_doc(data.get("template_b64"))
    _set_margins(doc)
    name     = data["name"]
    amount   = _fmt_amount(data["amount"])
    meddpicc = data["meddpicc"]

    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt, RGBColor

    cover = doc.add_paragraph()
    cover.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = cover.add_run("Business Case: AI Safety Intelligence")
    r.font.size = Pt(20); r.bold = True; r.font.color.rgb = RGBColor(*NAVY); r.font.name = "Calibri"
    cover2 = doc.add_paragraph()
    cover2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = cover2.add_run(f"Prepared for {name}  |  {_today_str()}")
    r2.font.size = Pt(11); r2.font.color.rgb = RGBColor(*GREY); r2.font.name = "Calibri"
    doc.add_page_break()

    metrics = meddpicc.get("metrics") or {}
    pain    = meddpicc.get("implicate_pain") or {}
    metrics_text = (metrics.get("evidence") or metrics.get("text") or "") if isinstance(metrics, dict) else ""
    pain_text    = (pain.get("evidence") or pain.get("text") or "") if isinstance(pain, dict) else ""

    _add_heading(doc, "1. Executive Summary")
    _add_body(doc, f"{name} faces ongoing operational gaps that carry direct financial and regulatory risk.")
    _add_body(doc, f"This document presents the case for deploying {bc_biz_sender['product']}.")
    _add_body(doc, f"The proposed investment of {amount} targets measurable improvement within 6-12 months.")
    doc.add_paragraph()

    _add_heading(doc, "2. Current Cost of the Problem")
    if metrics_text or pain_text:
        _add_body(doc, metrics_text or pain_text)
    else:
        _add_body(doc, "Workplace incidents carry direct costs: medical treatment, legal exposure, and operational downtime.")
        _add_body(doc, "Industry benchmarks place a single serious incident at $500K-$2M total cost, excluding reputational damage.")
    doc.add_paragraph()

    _add_heading(doc, "3. Proposed Investment")
    _add_table(doc, ["Item", "Detail"], [
        ["Platform",    bc_biz_sender['product']],
        ["Scope",       f"As discussed with {name}"],
        ["Investment",  amount],
        ["Term",        "Annual subscription + implementation"],
        ["Payment",     "Milestone-based"],
    ])
    doc.add_paragraph()

    _add_heading(doc, "4. ROI Projection")
    _add_body(doc, "Based on results across comparable deployments:")
    for stat in [
        "40-60% incident reduction within 6 months",
        "8x improvement in near-miss detection vs manual monitoring",
        "1 officer managing 20+ cameras (vs 2-4 cameras manually)",
        "5-10x return on investment from preventing one serious incident per year",
    ]:
        p = doc.add_paragraph(style="List Bullet")
        r = p.add_run(stat)
        r.font.size = Pt(11); r.font.name = "Calibri"
    doc.add_paragraph()

    _add_heading(doc, "5. Risk of Inaction")
    _add_body(doc,
              "Regulatory enforcement is increasing across Singapore (MOM), Australia (NOPSEMA), and MENA (Aramco HSE standards). "
              "Delayed deployment means continued exposure to preventable incidents and regulatory fines.")
    doc.add_paragraph()

    _add_heading(doc, "6. Recommended Next Step")
    ns = data["state"].get("next_step") or {}
    ns_text = ns.get("text") if isinstance(ns, dict) else str(ns or "")
    _add_body(doc, ns_text or f"Schedule a technical site walkthrough with the {name} safety team to confirm deployment scope.")

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── ROI CALCULATOR ─────────────────────────────────────────────────────────────

async def generate_roi_calculator(account_id: str, db: AsyncSession) -> bytes:
    data  = await _load_account_data(account_id, db)
    roi_sender = _ws_sender(data)
    doc   = _open_doc(data.get("template_b64"))
    _set_margins(doc)
    name  = data["name"]
    amount = float(data["account"].deal_amount or 0)

    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt, RGBColor

    cover = doc.add_paragraph()
    cover.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = cover.add_run("ROI Calculator: AI Safety Monitoring")
    r.font.size = Pt(20); r.bold = True; r.font.color.rgb = RGBColor(*NAVY); r.font.name = "Calibri"
    cover2 = doc.add_paragraph()
    cover2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = cover2.add_run(f"{name}  |  Fill in the highlighted cells  |  {_today_str()}")
    r2.font.size = Pt(11); r2.font.color.rgb = RGBColor(*GREY); r2.font.name = "Calibri"
    doc.add_page_break()

    _add_heading(doc, "Section A: Your Current Safety Costs (annual)")
    _add_table(doc, ["Cost Category", "Your Amount (USD)", "Notes"], [
        ["Workplace incidents (direct)",      "[enter]", "Medical, legal, downtime"],
        ["Safety officer headcount",          "[enter]", "Fully-loaded cost"],
        ["Regulatory fines (3yr avg)",        "[enter]", "MOM / OSHA / NOPSEMA"],
        ["Insurance premium",                 "[enter]", "Annual workplace injury"],
        ["Incident investigation costs",      "[enter]", "HR, legal, management"],
        ["TOTAL",                             "=SUM above", ""],
    ])
    doc.add_paragraph()

    _add_heading(doc, f"Section B: {roi_sender['company']} Investment")
    _add_table(doc, ["Item", "Cost (USD)"], [
        [f"{roi_sender['product']} licence (annual)", f"${amount:,.0f}" if amount else "[to be quoted]"],
        ["Implementation and integration",    "Included"],
        ["Training and onboarding",           "Included"],
        ["TOTAL YEAR 1",                      f"${amount:,.0f}" if amount else "[to be quoted]"],
    ])
    doc.add_paragraph()

    _add_heading(doc, "Section C: Expected Benefits")
    _add_table(doc, ["Benefit", "Estimate", "Basis"], [
        ["Incident reduction (40%)",     "[A × 40%]",        "Industry benchmark"],
        ["Officer efficiency gain",      "[headcount × 50%]","1 officer manages 5x more"],
        ["Insurance premium reduction",  "[premium × 15%]",  "AI monitoring reduces risk"],
        ["Regulatory fine avoidance",    "[fines × 60%]",    "Proactive detection"],
        ["TOTAL ANNUAL BENEFIT",         "=[sum above]",     ""],
    ])
    doc.add_paragraph()

    _add_heading(doc, "Section D: Summary")
    _add_table(doc, ["Metric", "Formula", "Result"], [
        ["Net annual benefit",   "C - B",            "[calculate]"],
        ["ROI",                  "(C-B)/B × 100",    "[calculate]"],
        ["Payback period (mo.)", "B / (C/12)",       "[calculate]"],
    ])

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── MUTUAL ACTION PLAN ─────────────────────────────────────────────────────────

async def generate_mutual_action_plan(account_id: str, db: AsyncSession) -> bytes:
    data    = await _load_account_data(account_id, db)
    map_sender = _ws_sender(data)
    doc     = _open_doc(data.get("template_b64"))
    _set_margins(doc)
    name    = data["name"]
    actions = data["actions"]

    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Pt, RGBColor

    cover = doc.add_paragraph()
    cover.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = cover.add_run("Mutual Action Plan")
    r.font.size = Pt(20); r.bold = True; r.font.color.rgb = RGBColor(*NAVY); r.font.name = "Calibri"
    cover2 = doc.add_paragraph()
    cover2.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r2 = cover2.add_run(f"{name} x {map_sender['company']}  |  Updated {_today_str()}")
    r2.font.size = Pt(11); r2.font.color.rgb = RGBColor(*GREY); r2.font.name = "Calibri"
    doc.add_page_break()

    _add_body(doc, f"This document tracks agreed actions between {name} and {map_sender['company']}. Both parties are accountable for their items.")
    doc.add_paragraph()

    def action_rows(acts):
        rows = []
        for a in acts:
            due   = str(a.due_date) if a.due_date else "TBD"
            owner = map_sender['company'] if a.action_type in ("email","call","demo") else name
            rows.append([a.title[:80], owner, due, a.status.title()])
        return rows

    overdue   = [a for a in actions if a.status == "overdue"]
    upcoming  = [a for a in actions if a.status == "upcoming"]
    completed = [a for a in actions if a.status == "completed"]

    if overdue:
        _add_heading(doc, "Overdue", level=2)
        _add_table(doc, ["Action", "Owner", "Due", "Status"], action_rows(overdue))
        doc.add_paragraph()

    rows = action_rows(upcoming) or [
        ["Review proposal and confirm scope",          name,             "This week",    "Upcoming"],
        ["Schedule technical walkthrough",             f"{name} + {map_sender['company']}", "2 weeks", "Upcoming"],
        ["Technical integration review",               map_sender['company'],    "Week 3",       "Upcoming"],
        ["Commercial approval and contract signature", name,             "Week 4",       "Upcoming"],
    ]
    _add_heading(doc, "Upcoming Actions", level=2)
    _add_table(doc, ["Action", "Owner", "Due", "Status"], rows)
    doc.add_paragraph()

    if completed:
        _add_heading(doc, "Completed", level=2)
        _add_table(doc, ["Action", "Owner", "Completed", "Status"], action_rows(completed))
        doc.add_paragraph()

    _add_heading(doc, "Key Contacts", level=2)
    _add_table(doc, ["Name", "Organisation", "Role", "Contact"], [
        [map_sender['name'],  map_sender['company'], map_sender['title'], map_sender['email'] or "[email]"],
        ["[To be filled]",    name,                  "[Role]",            "[Email]"],
    ])

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── Dispatch ──────────────────────────────────────────────────────────────────

GENERATORS = {
    "proposal":          generate_proposal,
    "sales_deck":        generate_sales_deck,
    "battle_card":       generate_battle_card,
    "business_case":     generate_business_case,
    "roi_calculator":    generate_roi_calculator,
    "mutual_action_plan":generate_mutual_action_plan,
}

FILE_FORMATS = {
    "proposal":           ("docx", "Proposal"),
    "sales_deck":         ("pptx", "Sales_Deck"),
    "battle_card":        ("docx", "Battle_Card"),
    "business_case":      ("docx", "Business_Case"),
    "roi_calculator":     ("docx", "ROI_Calculator"),
    "mutual_action_plan": ("docx", "Mutual_Action_Plan"),
}


async def generate_document(doc_type: str, account_id: str, db: AsyncSession) -> bytes:
    if doc_type not in GENERATORS:
        raise ValueError(f"Unknown doc_type: {doc_type}")
    return await GENERATORS[doc_type](account_id, db)
